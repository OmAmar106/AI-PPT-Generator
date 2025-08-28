from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from io import BytesIO
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE

def pptx_to_markdown(pptx_bytes, output_dir="extracted_images"):
    prs = Presentation(BytesIO(pptx_bytes))
    os.makedirs(output_dir, exist_ok=True)
    markdown_lines = []

    for i, slide in enumerate(prs.slides, start=1):
        markdown_lines.append(f"# Slide {i}\n")
        
        if slide.background and slide.background.fill.type == MSO_FILL.PICTURE:
            fill = slide.background.fill
            bg_image = fill._blip.fillFormat.blip.blob
            bg_filename = f"{output_dir}/slide{i}.jpeg"
            with open(bg_filename, "wb") as f:
                f.write(bg_image)
            markdown_lines.append(f"![Background]({os.path.basename(bg_filename)})\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                markdown_lines.append(shape.text.strip() + "\n")
        
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                img_filename = f"{output_dir}/slide{i}_{shape.shape_id}.{ext}"
                os.makedirs(output_dir, exist_ok=True)
                with open(img_filename, "wb") as f:
                    f.write(image.blob)
                markdown_lines.append(f"![Image]({os.path.basename(img_filename)})\n")
        markdown_lines.append("\n---\n")
    print(markdown_lines)
    return "\n".join(markdown_lines)
