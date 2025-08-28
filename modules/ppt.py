import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from modules.llm import ask_llm
from PIL import Image

def create_ppt(slides_md_list, api_key, theme_color=(245, 245, 245)):
    theme_color = RGBColor(*theme_color)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_md in ask_llm(slides_md_list, api_key):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg_shape = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=0, top=0,
            width=prs.slide_width, height=prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = theme_color
        bg_shape.line.fill.background()

        lines = slide_md.strip().split("\n")
        if not lines:
            continue

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_tf = title_shape.text_frame
        title_text = lines[0].lstrip("# ").strip()
        title_tf.text = title_text
        title_tf.paragraphs[0].font.size = Pt(32)
        title_tf.paragraphs[0].font.bold = True

        text_lines = []
        img_path = None
        for line in lines[1:]:
            match = re.match(r'!\[.*?\]\((.*?)\)', line)
            if match and img_path is None:
                img_path = os.path.join("extracted_images", match.group(1))
            else:
                text_lines.append(line.lstrip("# ").strip())

        text_width = 5.5 if (img_path and os.path.exists(img_path)) else 9

        text_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(text_width), Inches(5))
        tf = text_shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        font_size = Pt(16)
        for line in text_lines:
            if not line:
                continue

            p = tf.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            p.font.size = font_size

            if line.startswith("- "):
                p.text = line[2:].strip()
                p.level = 0
                p.bullet = True
                p.space_after = Pt(6)
            else:
                runs = re.split(r'(\*\*.*?\*\*|\*.*?\*)', line)
                for run_text in runs:
                    run = p.add_run()
                    if run_text.startswith("**") and run_text.endswith("**"):
                        run.text = run_text[2:-2]
                        run.font.bold = True
                    elif run_text.startswith("*") and run_text.endswith("*"):
                        run.text = run_text[1:-1]
                        run.font.italic = True
                    else:
                        run.text = run_text
                p.space_after = Pt(6)

        if img_path and os.path.exists(img_path):
            with Image.open(img_path) as img:
                max_width_inch = 4
                max_height_inch = 5
                width_px = int(max_width_inch * 96)
                height_px = int((width_px / img.width) * img.height)
                if height_px > int(max_height_inch * 96):
                    height_px = int(max_height_inch * 96)
                    width_px = int((height_px / img.height) * img.width)
                img_resized_path = img_path + "_resized.png"
                img.resize((width_px, height_px), Image.Resampling.LANCZOS).save(img_resized_path)
                slide.shapes.add_picture(
                    img_resized_path,
                    Inches(6),
                    Inches(1.5),
                    width=Inches(width_px / 96),
                    height=Inches(height_px / 96)
                )

    return prs
